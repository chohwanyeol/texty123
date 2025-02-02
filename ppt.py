from pptx import Presentation

def pptx(ppt):
  prs = Presentation(ppt)
  result = []
  for slide in prs.slides:
    for shape in slide.shapes:
      if not shape.has_text_frame:
        continue
        for paragraph in shape.text_frame.paragraphs:
          result.append(paragraph.text)
        return result